"""
企业级多源知识文档加载器 (Enterprise Multi-Source Document Loader)

支持三大类数据源：
  1. 文件系统: PDF, Word, Excel, PPTX, CSV, JSON, JSONL, XML, Markdown, HTML, TXT, LOG, EML, MSG, PNG/JPG/TIFF(OCR)
  2. 关系数据库: MySQL, PostgreSQL, SQL Server, Oracle, SQLite（通过SQLAlchemy）
  3. REST API: HTTP/HTTPS接口，支持JSONPath提取、模板渲染、缓存

核心设计原则：
  - 所有数据源统一输出为 LangChain Document 对象
  - 通过扩展名/协议自动匹配加载器，无需手动指定类型
  - 向后兼容原有的 prepare_documents() 和 load_documents() 接口
  - 支持可选的 YAML 配置文件驱动（data_sources.yml）
"""

import json
import os
import re
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# ============================================================================
# 内置焊接知识文档（作为知识库兜底数据，即使没有外部文档也能回答基础问题）
# ============================================================================

BUILTIN_DOCS = [
    Document(
        page_content=(
            "焊接工艺参数推荐知识：Q235碳钢10mm板厚推荐采用CO2气体保护焊(GMAW)，"
            "焊接电流220-280A，电压28-32V，焊接速度300-500mm/min，"
            "焊丝ER70S-6直径1.2mm，保护气体CO2或80%Ar+20%CO2流量15-20L/min。"
            "一般不需要预热（环境温度>5℃时）。"
            "10mm板建议开V型坡口，钝边2mm，间隙2-3mm。多层焊时层间温度控制在150-200℃。"
            "Q235碳当量低，焊接性好。"
        ),
        metadata={"source": "内置知识库", "category": "工艺参数", "material": "Q235", "thickness": "10mm"},
    ),
    Document(
        page_content=(
            "焊接工艺参数推荐知识：Q235碳钢6mm板厚推荐采用CO2气体保护焊(GMAW)，"
            "焊接电流180-240A，电压24-28V，焊接速度400-600mm/min，"
            "焊丝ER70S-6直径1.2mm，保护气体CO2流量15-20L/min。"
            "不需要预热。6mm板可不开坡口，留1-2mm间隙，单面焊双面成型。"
        ),
        metadata={"source": "内置知识库", "category": "工艺参数", "material": "Q235", "thickness": "6mm"},
    ),
    Document(
        page_content=(
            "焊接工艺参数推荐知识：Q345低合金钢10mm板厚推荐采用CO2气体保护焊(GMAW)或埋弧自动焊，"
            "焊接电流240-300A，电压30-34V，焊接速度300-450mm/min，"
            "焊丝ER50-6直径1.2mm，保护气体80%Ar+20%CO2流量18-22L/min。"
            "环境温度<0℃时需预热100-150℃。"
            "Q345碳当量略高于Q235，属低合金高强钢，注意控制热输入避免热影响区脆化。"
        ),
        metadata={"source": "内置知识库", "category": "工艺参数", "material": "Q345", "thickness": "10mm"},
    ),
    Document(
        page_content=(
            "焊接工艺参数推荐知识：不锈钢SUS304板厚3mm推荐采用TIG焊(GTAW)，"
            "焊接电流90-130A，电压12-16V，焊接速度150-250mm/min，"
            "焊丝ER308直径1.6mm，保护气体纯Ar流量8-12L/min。"
            "不需要预热。不锈钢导热系数低，热膨胀系数大，注意控制变形。"
            "焊接时电流比碳钢小10-20%。"
        ),
        metadata={"source": "内置知识库", "category": "工艺参数", "material": "SUS304", "thickness": "3mm"},
    ),
    # ---- 缺陷诊断知识 ----
    Document(
        page_content=(
            "焊接缺陷诊断：气孔（Porosity）的常见原因包括——保护气体流量不足或过大、焊丝受潮、"
            "工件表面有油污或铁锈、风速过大吹散保护气。"
            "解决方案：调整气体流量至15-25L/min，使用前烘干焊丝（250-300℃烘干1-2小时），"
            "焊前清理工件表面去除油污铁锈，设置防风屏障。"
        ),
        metadata={"source": "内置知识库", "category": "缺陷诊断", "defect": "气孔"},
    ),
    Document(
        page_content=(
            "焊接缺陷诊断：夹渣（Slag Inclusion）的常见原因——前层焊道清理不干净、焊接电流过小、焊接速度过慢。"
            "解决方案：层间打磨清理、适当增大电流、提高焊接速度。"
        ),
        metadata={"source": "内置知识库", "category": "缺陷诊断", "defect": "夹渣"},
    ),
    Document(
        page_content=(
            "焊接缺陷诊断：未焊透（Incomplete Penetration）的常见原因——坡口角度过小、钝边过大、"
            "电流过小、焊接速度过快。"
            "解决方案：增大坡口角度至55-65°、减小钝边至1-2mm、增大电流10-20A、降低焊速。"
        ),
        metadata={"source": "内置知识库", "category": "缺陷诊断", "defect": "未焊透"},
    ),
    Document(
        page_content=(
            "焊接缺陷诊断：裂纹（Crack）的常见原因——拘束应力过大、焊缝含氢量高、冷却速度过快、母材碳当量高。"
            "解决方案：预热降低冷却速度、使用低氢焊丝、合理安排焊接顺序减小应力、焊后热处理。"
        ),
        metadata={"source": "内置知识库", "category": "缺陷诊断", "defect": "裂纹"},
    ),
    Document(
        page_content=(
            "焊接缺陷诊断：咬边（Undercut）的常见原因——电流过大、电弧过长、焊接速度过快、焊枪角度不当。"
            "解决方案：减小电流、压低电弧长度、降低焊速、调整焊枪角度至70-80°。"
        ),
        metadata={"source": "内置知识库", "category": "缺陷诊断", "defect": "咬边"},
    ),
    # ---- 焊接标准规范 ----
    Document(
        page_content=(
            "焊接标准规范：GB/T 985.1-2008《气焊、焊条电弧焊、气体保护焊和高能束焊的推荐坡口》；"
            "GB/T 19867.1-2005《电弧焊焊接工艺规程》；"
            "GB/T 3323-2005《金属熔化焊焊接接头射线照相》；"
            "ISO 9606-1:2017《焊工考试 熔化焊 第1部分:钢》；"
            "ASME BPVC Section IX《焊接、钎焊和粘接评定》。"
        ),
        metadata={"source": "内置知识库", "category": "标准规范"},
    ),
]


# ============================================================================
# FileDocumentLoader —— 多格式文件系统加载器
# ============================================================================

class FileDocumentLoader:
    """支持20+文件格式的文档加载器。

    通过扩展名映射表自动匹配对应的加载方法，新增格式只需添加映射条目。
    所有加载方法统一返回 List[Document]。
    """

    def __init__(self):
        self._loaders: Dict[str, Callable] = self._build_loader_map()

    def _build_loader_map(self) -> Dict[str, Callable]:
        """构建扩展名 → 加载方法的映射表"""
        return {
            ".pdf":  self._load_pdf,
            ".docx": self._load_docx,
            ".xlsx": self._load_xlsx,
            ".xls":  self._load_xlsx,
            ".pptx": self._load_pptx,
            ".txt":  self._load_text,
            ".csv":  self._load_csv,
            ".json": self._load_json,
            ".jsonl":self._load_jsonl,
            ".xml":  self._load_xml,
            ".md":   self._load_markdown,
            ".html": self._load_html,
            ".htm":  self._load_html,
            ".log":  self._load_text,
            ".eml":  self._load_email,
            ".msg":  self._load_msg,
            ".png":  self._load_image_ocr,
            ".jpg":  self._load_image_ocr,
            ".jpeg": self._load_image_ocr,
            ".tiff": self._load_image_ocr,
            ".tif":  self._load_image_ocr,
        }

    # ------------------------------------------------------------------
    # 公共接口
    # ------------------------------------------------------------------

    def load_file(self, file_path: str) -> List[Document]:
        """加载单个文件，自动根据扩展名选择加载器"""
        ext = Path(file_path).suffix.lower()
        loader = self._loaders.get(ext)
        if loader is None:
            raise ValueError(f"不支持的文件格式: {ext} (文件: {file_path})")
        return loader(file_path)

    def load_directory(self, dir_path: str, recursive: bool = True) -> List[Document]:
        """递归扫描目录，加载所有支持格式的文件

        Args:
            dir_path: 目录路径
            recursive: 是否递归子目录

        Returns:
            所有文档的 Document 列表
        """
        if not os.path.isdir(dir_path):
            raise NotADirectoryError(f"目录不存在: {dir_path}")

        documents: List[Document] = []
        supported_exts = set(self._loaders.keys())

        for root, _, files in os.walk(dir_path):
            for file in files:
                ext = Path(file).suffix.lower()
                if ext not in supported_exts:
                    continue
                file_path = os.path.join(root, file)
                try:
                    docs = self._loaders[ext](file_path)
                    documents.extend(docs)
                except Exception as e:
                    # 单个文件加载失败不影响其他文件
                    print(f"[WARN] 加载文件失败 {file_path}: {e}")

            if not recursive:
                break

        return documents

    @property
    def supported_extensions(self) -> List[str]:
        """返回支持的文件扩展名列表"""
        return sorted(self._loaders.keys())

    # ------------------------------------------------------------------
    # PDF 加载
    # ------------------------------------------------------------------

    def _load_pdf(self, file_path: str) -> List[Document]:
        """加载PDF文档（优先 PyMuPDF，回退 PyPDF）"""
        documents = []
        file_name = Path(file_path).name

        # 优先使用 PyMuPDF（速度更快、中文支持更好）
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                if text.strip():
                    documents.append(Document(
                        page_content=text.strip(),
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": "文档",
                            "page": page_num + 1,
                        }
                    ))
            doc.close()
            return documents
        except ImportError:
            pass
        except Exception as e:
            print(f"[WARN] PyMuPDF加载失败，尝试PyPDF: {e}")

        # 回退 PyPDF
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    documents.append(Document(
                        page_content=text.strip(),
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": "文档",
                            "page": page_num + 1,
                        }
                    ))
        except Exception as e:
            raise RuntimeError(f"PDF加载失败 {file_path}: {e}")

        return documents

    # ------------------------------------------------------------------
    # Word 加载
    # ------------------------------------------------------------------

    def _load_docx(self, file_path: str) -> List[Document]:
        """加载Word文档(.docx)"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            full_text = "\n".join(
                para.text for para in doc.paragraphs if para.text.strip()
            )
            if not full_text.strip():
                return []
            return [Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "file_name": Path(file_path).name,
                    "data_type": "file",
                    "category": "文档",
                }
            )]
        except Exception as e:
            raise RuntimeError(f"Word文档加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # Excel 加载
    # ------------------------------------------------------------------

    def _load_xlsx(self, file_path: str) -> List[Document]:
        """加载Excel文档(.xlsx/.xls)，每张表转为独立Document"""
        documents = []
        ext = Path(file_path).suffix.lower()
        file_name = Path(file_path).name

        try:
            if ext == ".xls":
                import pandas as pd
                excel_file = pd.ExcelFile(file_path, engine="xlrd")
            else:
                import pandas as pd
                excel_file = pd.ExcelFile(file_path, engine="openpyxl")

            for sheet_name in excel_file.sheet_names:
                df = excel_file.parse(sheet_name)
                if df.empty:
                    continue
                # 将DataFrame转为Markdown表格格式的文本（LLM友好）
                header = "| " + " | ".join(str(c) for c in df.columns) + " |"
                separator = "| " + " | ".join("---" for _ in df.columns) + " |"
                rows = []
                for _, row in df.iterrows():
                    rows.append("| " + " | ".join(
                        str(v) if pd.notna(v) else "" for v in row
                    ) + " |")
                table_text = f"[表: {sheet_name}]\n{header}\n{separator}\n" + "\n".join(rows)

                documents.append(Document(
                    page_content=table_text,
                    metadata={
                        "source": file_path,
                        "file_name": file_name,
                        "data_type": "file",
                        "category": "数据表",
                        "sheet": sheet_name,
                        "rows": len(df),
                    }
                ))

        except ImportError as e:
            raise ImportError(f"Excel加载需要pandas和openpyxl/xlrd: {e}")
        except Exception as e:
            raise RuntimeError(f"Excel文档加载失败 {file_path}: {e}")

        return documents

    # ------------------------------------------------------------------
    # PowerPoint 加载
    # ------------------------------------------------------------------

    def _load_pptx(self, file_path: str) -> List[Document]:
        """加载PowerPoint文档(.pptx)，按幻灯片提取文本"""
        documents = []
        file_name = Path(file_path).name

        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    documents.append(Document(
                        page_content=f"[幻灯片 {slide_num + 1}]\n" + "\n".join(texts),
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": "演示文档",
                            "slide": slide_num + 1,
                        }
                    ))
        except ImportError:
            raise ImportError("PPTX加载需要 python-pptx 库: pip install python-pptx")
        except Exception as e:
            raise RuntimeError(f"PPTX加载失败 {file_path}: {e}")

        return documents

    # ------------------------------------------------------------------
    # 纯文本 / Markdown / 日志
    # ------------------------------------------------------------------

    def _load_text(self, file_path: str) -> List[Document]:
        """加载纯文本文件(.txt/.log)"""
        file_name = Path(file_path).name
        ext = Path(file_path).suffix.lower()
        category = "日志" if ext == ".log" else "纯文本"

        # 尝试多种编码
        for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                if text.strip():
                    return [Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": category,
                            "encoding": encoding,
                        }
                    )]
                return []
            except UnicodeDecodeError:
                continue

        raise RuntimeError(f"无法解码文件 {file_path}，尝试了 utf-8/gbk/gb2312/latin-1 均失败")

    def _load_markdown(self, file_path: str) -> List[Document]:
        """加载Markdown文件，保留原始格式"""
        file_name = Path(file_path).name
        for encoding in ["utf-8", "gbk", "gb2312"]:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                if text.strip():
                    return [Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": "Markdown文档",
                        }
                    )]
                return []
            except UnicodeDecodeError:
                continue
        raise RuntimeError(f"无法解码Markdown文件 {file_path}")

    # ------------------------------------------------------------------
    # CSV
    # ------------------------------------------------------------------

    def _load_csv(self, file_path: str) -> List[Document]:
        """加载CSV文件"""
        file_name = Path(file_path).name
        try:
            import pandas as pd
            # 尝试多种编码
            for encoding in ["utf-8", "gbk", "gb2312"]:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    break
                except (UnicodeDecodeError, Exception):
                    continue
            else:
                df = pd.read_csv(file_path, encoding="latin-1")

            if df.empty:
                return []

            # 转为可读文本
            header = "| " + " | ".join(str(c) for c in df.columns) + " |"
            separator = "| " + " | ".join("---" for _ in df.columns) + " |"
            rows = []
            for _, row in df.iterrows():
                rows.append("| " + " | ".join(
                    str(v) if pd.notna(v) else "" for v in row
                ) + " |")
            text = f"[CSV文件: {file_name}]\n{header}\n{separator}\n" + "\n".join(rows)

            return [Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "file_name": file_name,
                    "data_type": "file",
                    "category": "数据表",
                    "rows": len(df),
                }
            )]
        except ImportError:
            raise ImportError("CSV加载需要 pandas 库: pip install pandas")
        except Exception as e:
            raise RuntimeError(f"CSV加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # JSON / JSONL
    # ------------------------------------------------------------------

    def _load_json(self, file_path: str) -> List[Document]:
        """加载JSON文件，支持对象和数组两种结构"""
        file_name = Path(file_path).name
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            documents = []
            if isinstance(data, list):
                # 数组：每个元素一个Document
                for idx, item in enumerate(data):
                    if isinstance(item, dict):
                        content = json.dumps(item, ensure_ascii=False, indent=2)
                    else:
                        content = str(item)
                    documents.append(Document(
                        page_content=content,
                        metadata={
                            "source": file_path,
                            "file_name": file_name,
                            "data_type": "file",
                            "category": "JSON数据",
                            "index": idx,
                        }
                    ))
            elif isinstance(data, dict):
                # 对象：整个为一个Document
                documents.append(Document(
                    page_content=json.dumps(data, ensure_ascii=False, indent=2),
                    metadata={
                        "source": file_path,
                        "file_name": file_name,
                        "data_type": "file",
                        "category": "JSON数据",
                    }
                ))

            return documents
        except Exception as e:
            raise RuntimeError(f"JSON加载失败 {file_path}: {e}")

    def _load_jsonl(self, file_path: str) -> List[Document]:
        """加载JSONL文件（每行一个JSON对象）"""
        file_name = Path(file_path).name
        documents = []
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                for line_num, line in enumerate(f):
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        obj = json.loads(line)
                        if isinstance(obj, dict):
                            content = json.dumps(obj, ensure_ascii=False, indent=2)
                        else:
                            content = str(obj)
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                "source": file_path,
                                "file_name": file_name,
                                "data_type": "file",
                                "category": "JSONL数据",
                                "line": line_num + 1,
                            }
                        ))
                    except json.JSONDecodeError:
                        continue
            return documents
        except Exception as e:
            raise RuntimeError(f"JSONL加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # XML
    # ------------------------------------------------------------------

    def _load_xml(self, file_path: str) -> List[Document]:
        """加载XML文件，递归提取文本内容"""
        file_name = Path(file_path).name
        try:
            # 优先使用 lxml（更快）
            try:
                import lxml.etree as ET
                tree = ET.parse(file_path)
                root = tree.getroot()
            except ImportError:
                import xml.etree.ElementTree as ET
                tree = ET.parse(file_path)
                root = tree.getroot()

            # 递归提取所有文本
            def extract_text(element, depth=0):
                texts = []
                if element.text and element.text.strip():
                    texts.append(f"{'  ' * depth}{element.tag}: {element.text.strip()}")
                for child in element:
                    texts.extend(extract_text(child, depth + 1))
                return texts

            text_lines = extract_text(root)
            if not text_lines:
                return []
            return [Document(
                page_content=f"[XML文件: {file_name}]\n" + "\n".join(text_lines),
                metadata={
                    "source": file_path,
                    "file_name": file_name,
                    "data_type": "file",
                    "category": "XML数据",
                }
            )]
        except Exception as e:
            raise RuntimeError(f"XML加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # HTML
    # ------------------------------------------------------------------

    def _load_html(self, file_path: str) -> List[Document]:
        """加载HTML文件，去除标签提取纯文本"""
        file_name = Path(file_path).name
        try:
            from bs4 import BeautifulSoup
            for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        html = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise RuntimeError(f"无法解码HTML文件 {file_path}")

            soup = BeautifulSoup(html, "lxml" if _has_lxml() else "html.parser")
            # 移除script和style标签
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            # 清理多余空行
            lines = [line.strip() for line in text.split("\n") if line.strip()]
            clean_text = "\n".join(lines)

            if not clean_text:
                return []
            return [Document(
                page_content=clean_text,
                metadata={
                    "source": file_path,
                    "file_name": file_name,
                    "data_type": "file",
                    "category": "网页文档",
                }
            )]
        except ImportError:
            raise ImportError("HTML加载需要 beautifulsoup4 和 lxml: pip install beautifulsoup4 lxml")
        except Exception as e:
            raise RuntimeError(f"HTML加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # 邮件 (EML / MSG)
    # ------------------------------------------------------------------

    def _load_email(self, file_path: str) -> List[Document]:
        """加载EML邮件文件"""
        import email
        from email import policy
        file_name = Path(file_path).name

        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                msg = email.message_from_string(f.read(), policy=policy.default)

            parts = [f"主题: {msg.get('Subject', '(无主题)')}"]
            parts.append(f"发件人: {msg.get('From', '')}")
            parts.append(f"收件人: {msg.get('To', '')}")
            parts.append(f"日期: {msg.get('Date', '')}")

            # 提取正文
            body_parts = []
            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        try:
                            payload = part.get_payload(decode=True)
                            charset = part.get_content_charset() or "utf-8"
                            body_parts.append(payload.decode(charset, errors="replace"))
                        except Exception:
                            pass
            else:
                try:
                    payload = msg.get_payload(decode=True)
                    charset = msg.get_content_charset() or "utf-8"
                    body_parts.append(payload.decode(charset, errors="replace"))
                except Exception:
                    pass

            if body_parts:
                parts.append("\n--- 正文 ---\n" + "\n".join(body_parts))

            content = "\n".join(parts)
            if not content.strip():
                return []
            return [Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "file_name": file_name,
                    "data_type": "file",
                    "category": "邮件",
                }
            )]
        except Exception as e:
            raise RuntimeError(f"EML邮件加载失败 {file_path}: {e}")

    def _load_msg(self, file_path: str) -> List[Document]:
        """加载MSG邮件文件（Outlook格式）"""
        file_name = Path(file_path).name
        try:
            import extract_msg
            msg = extract_msg.Message(file_path)
            parts = [
                f"主题: {msg.subject or '(无主题)'}",
                f"发件人: {msg.sender}",
                f"收件人: {msg.to}",
                f"日期: {msg.date}",
            ]
            if msg.body:
                parts.append(f"\n--- 正文 ---\n{msg.body}")
            content = "\n".join(parts)
            return [Document(
                page_content=content,
                metadata={
                    "source": file_path,
                    "file_name": file_name,
                    "data_type": "file",
                    "category": "邮件",
                }
            )]
        except ImportError:
            raise ImportError("MSG加载需要 extract-msg: pip install extract-msg")
        except Exception as e:
            raise RuntimeError(f"MSG邮件加载失败 {file_path}: {e}")

    # ------------------------------------------------------------------
    # 图片 OCR
    # ------------------------------------------------------------------

    def _load_image_ocr(self, file_path: str) -> List[Document]:
        """通过OCR从图片提取文字"""
        file_name = Path(file_path).name
        try:
            from PIL import Image
            img = Image.open(file_path)
            # 尝试导入pytesseract
            try:
                import pytesseract
                text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            except ImportError:
                raise ImportError(
                    "图片OCR需要 pytesseract 和 Pillow，以及系统安装Tesseract-OCR:\n"
                    "  pip install pytesseract Pillow\n"
                    "  Windows: https://github.com/UB-Mannheim/tesseract/wiki"
                )
            if text.strip():
                return [Document(
                    page_content=text.strip(),
                    metadata={
                        "source": file_path,
                        "file_name": file_name,
                        "data_type": "image_ocr",
                        "category": "图片识别",
                    }
                )]
            return []
        except ImportError:
            raise ImportError("图片OCR需要 Pillow: pip install Pillow")
        except Exception as e:
            raise RuntimeError(f"图片OCR处理失败 {file_path}: {e}")


# ============================================================================
# DatabaseDocumentLoader —— 关系数据库数据抽取
# ============================================================================

class DatabaseDocumentLoader:
    """从关系数据库加载知识数据。

    支持通过SQLAlchemy连接：MySQL, PostgreSQL, SQL Server, Oracle, SQLite 等。
    提供两种抽取模式：
      - 表导出模式: load_table()  将整张表转为Document列表
      - SQL查询模式: load_query()  通过自定义SQL精确提取
    """

    def __init__(self):
        self.engine = None

    def connect(self, db_url: str):
        """建立数据库连接

        Args:
            db_url: SQLAlchemy连接字符串
                    示例:
                      - SQLite:     sqlite:///data/welding_agent.db
                      - MySQL:      mysql+pymysql://user:pass@host:3306/db
                      - PostgreSQL: postgresql+psycopg2://user:pass@host:5432/db
                      - SQL Server: mssql+pyodbc://user:pass@host:1433/db?driver=ODBC+Driver+17
                      - Oracle:     oracle+cx_oracle://user:pass@host:1521/service_name
        """
        try:
            from sqlalchemy import create_engine
            self.engine = create_engine(db_url, connect_args={"check_same_thread": False}
                                        if "sqlite" in db_url else {})
        except ImportError:
            raise ImportError("数据库加载需要 SQLAlchemy: pip install sqlalchemy")

    def load_table(
        self,
        table_name: str,
        description: str = "",
        limit: int = 10000,
        columns: Optional[List[str]] = None,
        where: Optional[str] = None,
    ) -> List[Document]:
        """加载数据库表，每行转为一条Document

        Args:
            table_name: 表名
            description: 数据描述（写入metadata）
            limit: 最大加载行数（防止超大表OOM）
            columns: 指定加载的列（None=全部列）
            where: SQL WHERE 条件（不含WHERE关键字）

        Returns:
            Document列表，每行一个，内容为 "列名=值, ..." 格式
        """
        if self.engine is None:
            raise RuntimeError("未连接数据库，请先调用 connect()")

        try:
            import pandas as pd
            cols = ", ".join(columns) if columns else "*"
            sql = f"SELECT {cols} FROM {table_name}"
            if where:
                sql += f" WHERE {where}"
            if limit:
                sql += f" LIMIT {limit}"

            df = pd.read_sql(sql, self.engine)
            if df.empty:
                return []

            # 如果指定了limit且超过限制，则采样
            if limit and len(df) > limit:
                df = df.sample(n=limit, random_state=42)

            documents = []
            for idx, row in df.iterrows():
                parts = [f"[{table_name}]"]
                for col, val in row.items():
                    if pd.notna(val):
                        parts.append(f"{col}={val}")
                content = ", ".join(parts)
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "source": f"db://{table_name}",
                        "data_type": "database",
                        "category": description or table_name,
                        "row_id": int(idx) if hasattr(idx, "__int__") else idx,
                        "table": table_name,
                    }
                ))
            return documents

        except ImportError:
            raise ImportError("数据库加载需要 pandas: pip install pandas")
        except Exception as e:
            raise RuntimeError(f"数据库表加载失败 [{table_name}]: {e}")

    def load_query(self, sql: str, description: str = "") -> List[Document]:
        """执行自定义SQL查询并转为Document

        Args:
            sql: SQL查询语句
            description: 数据描述

        Returns:
            Document列表
        """
        if self.engine is None:
            raise RuntimeError("未连接数据库，请先调用 connect()")

        try:
            import pandas as pd
            df = pd.read_sql(sql, self.engine)
            if df.empty:
                return []

            documents = []
            for idx, row in df.iterrows():
                parts = []
                for col, val in row.items():
                    if pd.notna(val):
                        parts.append(f"{col}={val}")
                content = ", ".join(parts)
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "source": "db://custom_query",
                        "data_type": "database",
                        "category": description or "SQL查询结果",
                        "row_id": int(idx) if hasattr(idx, "__int__") else idx,
                    }
                ))
            return documents

        except ImportError:
            raise ImportError("数据库加载需要 pandas: pip install pandas")
        except Exception as e:
            raise RuntimeError(f"SQL查询执行失败: {e}")


# ============================================================================
# APIDocumentLoader —— REST API 数据接入
# ============================================================================

class APIDocumentLoader:
    """通过REST API加载知识数据。

    支持功能：
      - JSONPath 响应路径提取
      - Jinja2 模板渲染格式化
      - 内存缓存（避免重复请求）
      - 自定义请求头（认证等）
    """

    def __init__(self):
        self._cache: Dict[str, List[Document]] = {}

    def load_from_api(
        self,
        url: str,
        method: str = "GET",
        headers: Optional[Dict[str, str]] = None,
        params: Optional[Dict[str, Any]] = None,
        json_body: Optional[Dict[str, Any]] = None,
        json_path: str = "$",
        template: Optional[str] = None,
        description: str = "",
        cache_key: Optional[str] = None,
        timeout: int = 30,
    ) -> List[Document]:
        """调用REST API并解析响应为Document列表

        Args:
            url: API地址（支持环境变量：${VAR_NAME}）
            method: HTTP方法 (GET/POST/PUT)
            headers: 请求头字典
            params: URL查询参数
            json_body: POST请求的JSON body
            json_path: JSONPath表达式，用于从响应中提取数据
                        - "$"         : 整个响应
                        - "$.data[*]" : data数组的所有元素
                        - "$.results[0:10]": 前10条
            template: Jinja2模板字符串，格式化每条数据
                       - "{{ item.title }}: {{ item.content }}"
                       - 不指定时序列化为JSON字符串
            description: 数据描述（写入metadata）
            cache_key: 缓存键，相同key跳过请求直接返回缓存
            timeout: 请求超时秒数

        Returns:
            Document列表

        使用示例:
            loader = APIDocumentLoader()
            docs = loader.load_from_api(
                url="http://internal-api/welding/standards",
                json_path="$.data[*]",
                template="标准 {{ item.code }}: {{ item.name }}",
                description="焊接标准规范",
                cache_key="welding_standards",
            )
        """
        try:
            import requests
            from jsonpath_ng import parse as jsonpath_parse
        except ImportError:
            raise ImportError("API加载需要 requests 和 jsonpath-ng: pip install requests jsonpath-ng")

        # 环境变量替换
        if "${" in url:
            url = _resolve_env_vars(url)
        if headers:
            headers = {k: _resolve_env_vars(v) if isinstance(v, str) else v
                       for k, v in headers.items()}

        # 缓存检查
        if cache_key and cache_key in self._cache:
            return self._cache[cache_key]

        # 发送请求
        try:
            resp = requests.request(
                method=method.upper(),
                url=url,
                headers=headers or {},
                params=params or {},
                json=json_body or {},
                timeout=timeout,
            )
            resp.raise_for_status()
            data = resp.json()
        except requests.exceptions.RequestException as e:
            raise RuntimeError(f"API请求失败 [{method} {url}]: {e}")

        # JSONPath提取
        try:
            matches = jsonpath_parse(json_path).find(data)
            items = [m.value for m in matches]
        except Exception as e:
            raise RuntimeError(f"JSONPath解析失败 [{json_path}]: {e}")

        if not items:
            return []

        # 如果只有一个元素且不是列表，包裹一下
        if not isinstance(items, list) or len(items) == 0:
            items = [items]

        # 转换为Document
        documents = []
        if template:
            try:
                from jinja2 import Template
                tmpl = Template(template)
                for i, item in enumerate(items):
                    content = tmpl.render(item=item, index=i)
                    if content.strip():
                        documents.append(Document(
                            page_content=content,
                            metadata={
                                "source": url,
                                "data_type": "api",
                                "category": description or "API数据",
                                "index": i,
                            }
                        ))
            except ImportError:
                raise ImportError("模板渲染需要 Jinja2: pip install Jinja2")
            except Exception as e:
                raise RuntimeError(f"模板渲染失败: {e}")
        else:
            for i, item in enumerate(items):
                if isinstance(item, dict):
                    content = json.dumps(item, ensure_ascii=False, indent=2)
                else:
                    content = str(item)
                documents.append(Document(
                    page_content=content,
                    metadata={
                        "source": url,
                        "data_type": "api",
                        "category": description or "API数据",
                        "index": i,
                    }
                ))

        # 更新缓存
        if cache_key:
            self._cache[cache_key] = documents

        return documents

    def clear_cache(self):
        """清空缓存"""
        self._cache.clear()


# ============================================================================
# UnifiedKnowledgeLoader —— 统一知识加载入口
# ============================================================================

class UnifiedKnowledgeLoader:
    """统一知识加载器——多源数据一站式入口。

    整合 文件系统 + 关系数据库 + REST API 三大数据源，
    自动识别、加载、统一输出为 LangChain Document 列表。
    """

    def __init__(self):
        self.file_loader = FileDocumentLoader()
        self.db_loader = DatabaseDocumentLoader()
        self.api_loader = APIDocumentLoader()

    def load_all(
        self,
        docs_dir: Optional[str] = None,
        db_configs: Optional[List[Dict[str, Any]]] = None,
        api_configs: Optional[List[Dict[str, Any]]] = None,
        include_builtin: bool = True,
    ) -> List[Document]:
        """从所有配置的数据源加载知识文档

        Args:
            docs_dir: 文件目录路径，为空则不加载文件
            db_configs: 数据库配置列表，每一项包含 url + tables/queries
            api_configs: API配置列表，每一项为 load_from_api() 的参数字典
            include_builtin: 是否包含内置焊接知识文档

        Returns:
            统一的 Document 列表

        典型调用:
            loader = UnifiedKnowledgeLoader()
            documents = loader.load_all(
                docs_dir="./data/knowledge_docs",
                db_configs=[
                    {
                        "url": "mysql+pymysql://reader:pass@192.168.1.100:3306/welding_mes",
                        "tables": [
                            {"name": "welding_records", "description": "焊接生产记录", "limit": 5000},
                        ],
                    },
                ],
                api_configs=[
                    {
                        "url": "http://internal-api/welding/standards",
                        "json_path": "$.data[*]",
                        "template": "标准 {{ item.code }}: {{ item.name }}",
                        "description": "焊接标准",
                    },
                ],
            )
        """
        all_docs: List[Document] = []

        # ---- 1. 内置知识 ----
        if include_builtin:
            all_docs.extend(BUILTIN_DOCS)

        # ---- 2. 文件加载 ----
        if docs_dir and os.path.isdir(docs_dir):
            try:
                file_docs = self.file_loader.load_directory(docs_dir)
                all_docs.extend(file_docs)
            except Exception as e:
                print(f"[WARN] 文件目录加载失败 {docs_dir}: {e}")

        # ---- 3. 数据库加载 ----
        if db_configs:
            for db_cfg in db_configs:
                url = db_cfg.get("url", "")
                if not url:
                    continue
                try:
                    # 环境变量替换
                    url = _resolve_env_vars(url)
                    self.db_loader.connect(url)

                    # 加载表
                    for table_cfg in db_cfg.get("tables", []):
                        table_name = table_cfg["name"]
                        docs = self.db_loader.load_table(
                            table_name=table_name,
                            description=table_cfg.get("description", ""),
                            limit=table_cfg.get("limit", 10000),
                            columns=table_cfg.get("columns"),
                            where=table_cfg.get("where"),
                        )
                        all_docs.extend(docs)

                    # 执行自定义查询
                    for query_cfg in db_cfg.get("queries", []):
                        sql = query_cfg["sql"]
                        docs = self.db_loader.load_query(
                            sql=sql,
                            description=query_cfg.get("description", ""),
                        )
                        all_docs.extend(docs)

                except Exception as e:
                    db_id = db_cfg.get("id", url)
                    print(f"[WARN] 数据库加载失败 [{db_id}]: {e}")

        # ---- 4. API加载 ----
        if api_configs:
            for api_cfg in api_configs:
                try:
                    url = api_cfg.get("url", "")
                    if not url:
                        continue
                    docs = self.api_loader.load_from_api(**api_cfg)
                    all_docs.extend(docs)
                except Exception as e:
                    api_id = api_cfg.get("id", api_cfg.get("url", ""))
                    print(f"[WARN] API加载失败 [{api_id}]: {e}")

        return all_docs


# ============================================================================
# 公共接口（向后兼容）
# ============================================================================

def load_document(file_path: str, source: Optional[str] = None) -> List[Document]:
    """加载单个文档文件（用于API上传场景）

    Args:
        file_path: 文件路径
        source: 来源标识（可选）

    Returns:
        Document列表
    """
    loader = FileDocumentLoader()
    documents = loader.load_file(file_path)

    # 如果指定了source，更新metadata
    if source:
        for doc in documents:
            doc.metadata["source"] = source

    return documents


def prepare_documents(docs_dir: str) -> List[Document]:
    """准备知识库文档（向后兼容旧接口）

    从指定目录加载所有支持的文档格式，并与内置知识库合并。

    Args:
        docs_dir: 知识文档目录

    Returns:
        Document列表
    """
    return prepare_multi_source_documents(docs_dir=docs_dir)


def load_documents(docs_dir: str) -> List[Document]:
    """加载知识库文档（向后兼容旧接口别名）"""
    return prepare_documents(docs_dir)


def prepare_multi_source_documents(
    docs_dir: Optional[str] = None,
    db_configs: Optional[List[Dict[str, Any]]] = None,
    api_configs: Optional[List[Dict[str, Any]]] = None,
    include_builtin: bool = True,
) -> List[Document]:
    """从多源加载知识文档（推荐的新接口）

    Args:
        docs_dir: 文件目录路径
        db_configs: 数据库配置列表
        api_configs: API配置列表
        include_builtin: 是否包含内置焊接知识

    Returns:
        统一的 Document 列表
    """
    loader = UnifiedKnowledgeLoader()
    return loader.load_all(
        docs_dir=docs_dir,
        db_configs=db_configs,
        api_configs=api_configs,
        include_builtin=include_builtin,
    )


# ============================================================================
# 文档分块
# ============================================================================

def split_documents(
    documents: List[Document],
    chunk_size: int = 500,
    chunk_overlap: int = 50,
) -> List[Document]:
    """将文档列表按指定大小分割成块，用于向量化

    Args:
        documents: 待分割的Document列表
        chunk_size: 每块最大字符数
        chunk_overlap: 相邻块重叠字符数

    Returns:
        分割后的Document列表，每个chunk的metadata会继承原文档metadata并附加chunk_id
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", "。", ".", "；", ";", " ", ""],
        length_function=len,
    )
    chunks = splitter.split_documents(documents)

    # 为每个chunk添加唯一ID
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_id"] = i

    return chunks


# ============================================================================
# 辅助函数
# ============================================================================

def _has_lxml() -> bool:
    """检查是否安装了lxml"""
    try:
        import lxml  # noqa
        return True
    except ImportError:
        return False


def _resolve_env_vars(value: str) -> str:
    """解析字符串中的环境变量占位符 ${VAR_NAME}"""
    pattern = re.compile(r'\$\{([^}]+)\}')
    def replacer(match):
        var_name = match.group(1)
        return os.environ.get(var_name, match.group(0))
    return pattern.sub(replacer, value)


def load_data_sources_config(config_path: str) -> Dict[str, Any]:
    """从YAML配置文件加载数据源配置

    配置文件格式见 data_sources.yml。

    Args:
        config_path: YAML配置文件路径

    Returns:
        解析后的配置字典，包含 docs_dir, databases, apis, chunking 等字段
    """
    try:
        import yaml
    except ImportError:
        # 如果没有PyYAML，尝试内置JSON方式
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except json.JSONDecodeError:
            raise ImportError("YAML配置需要 PyYAML: pip install PyYAML")

    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    # 处理环境变量
    return _resolve_config_env(config)


def _resolve_config_env(config: Any, depth: int = 0) -> Any:
    """递归解析配置中的环境变量"""
    if depth > 10:
        return config
    if isinstance(config, dict):
        return {k: _resolve_config_env(v, depth + 1) for k, v in config.items()}
    elif isinstance(config, list):
        return [_resolve_config_env(item, depth + 1) for item in config]
    elif isinstance(config, str):
        return _resolve_env_vars(config)
    return config


def prepare_from_config(config_path: str) -> List[Document]:
    """从YAML配置文件加载多源知识文档

    这是最简单的使用方式，只需一行代码：

        documents = prepare_from_config("data_sources.yml")

    Args:
        config_path: YAML配置文件路径

    Returns:
        Document列表
    """
    config = load_data_sources_config(config_path)

    docs_dir = config.get("docs_dir")
    db_configs = config.get("databases", [])
    api_configs = config.get("apis", [])

    # 过滤掉 disabled 的数据源
    db_configs = [db for db in db_configs if db.get("enabled", True)]
    api_configs = [api for api in api_configs if api.get("enabled", True)]

    return prepare_multi_source_documents(
        docs_dir=docs_dir,
        db_configs=db_configs,
        api_configs=api_configs,
        include_builtin=True,
    )
